#!/usr/bin/env python3
"""
Genera «Aplicación Monitoreo ReTo.pptx» — presentación profesional (5 slides).

Requisitos:  pip install python-pptx
Uso:         python3 generar_presentacion_reto.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = Path(__file__).resolve().parent
OUT_NAME = "Aplicación Monitoreo ReTo.pptx"
LOGO_RETO = SCRIPT_DIR / "logo_reto.png"
LOGOS_DIR = SCRIPT_DIR / "logos"

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(0x1A, 0x2C, 0x3E)
NAVY_LIGHT = RGBColor(0x22, 0x3A, 0x52)
BLUE = RGBColor(0x2D, 0x89, 0xC8)
ORANGE = RGBColor(0xF5, 0x82, 0x20)
GREEN = RGBColor(0x2E, 0x93, 0x48)
PURPLE = RGBColor(0x7D, 0x3C, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY_BG = RGBColor(0xF5, 0xF5, 0xF5)
GREY_TEXT = RGBColor(0x55, 0x55, 0x55)
DARK_TEXT = RGBColor(0x2A, 0x2A, 0x2A)
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xE0)


def _rect(slide, left, top, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def _rounded_rect(slide, left, top, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def _circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _text_box(slide, left, top, w, h):
    box = slide.shapes.add_textbox(left, top, w, h)
    box.text_frame.word_wrap = True
    return box.text_frame


def _add_text(tf, text, size, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, space_after=4):
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(space_after)
    return p


def _add_bullet(tf, text, size=16, color=DARK_TEXT, bold=False, bullet_char="\u2022"):
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{bullet_char}  {text}"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    p.space_after = Pt(5)
    p.space_before = Pt(2)
    return p


def _add_footer_logos(slide):
    """Agrega barra inferior con logos de socios."""
    _rect(slide, Inches(0), H - Inches(0.85), W, Inches(0.85), GREY_BG)
    _rect(slide, Inches(0), H - Inches(0.86), W, Inches(0.02), ORANGE)

    x = Inches(0.6)
    logo_h = Inches(0.5)

    if LOGOS_DIR.exists():
        for png in sorted(LOGOS_DIR.glob("*.png")):
            try:
                slide.shapes.add_picture(str(png), x, H - Inches(0.72), height=logo_h)
                x += Inches(1.6)
                if float(x) > float(W) - Inches(1):
                    break
            except Exception:
                pass


def _slide_title_bar(slide, title_text, subtitle_text=None):
    """Barra superior azul oscuro con título blanco."""
    _rect(slide, Inches(0), Inches(0), W, Inches(1.25), NAVY)
    _rect(slide, Inches(0), Inches(1.25), W, Inches(0.06), ORANGE)

    tf = _text_box(slide, Inches(0.8), Inches(0.18), Inches(11.5), Inches(0.7))
    _add_text(tf, title_text, 28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle_text:
        tf2 = _text_box(slide, Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.4))
        _add_text(tf2, subtitle_text, 14, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.LEFT)


def _step_card(slide, x, y, number, title, desc, accent_color):
    """Tarjeta de paso con número circular + texto."""
    card = _rounded_rect(slide, x, y, Inches(2.65), Inches(3.6), WHITE)
    card.shadow.inherit = False

    circ = _circle(slide, x + Inches(0.85), y + Inches(0.25), Inches(0.9), accent_color)
    ctf = circ.text_frame
    ctf.paragraphs[0].text = str(number)
    ctf.paragraphs[0].font.size = Pt(28)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.word_wrap = False

    tf = _text_box(slide, x + Inches(0.15), y + Inches(1.35), Inches(2.35), Inches(0.5))
    _add_text(tf, title, 18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    tf2 = _text_box(slide, x + Inches(0.15), y + Inches(1.85), Inches(2.35), Inches(1.5))
    _add_text(tf2, desc, 13, color=GREY_TEXT, align=PP_ALIGN.CENTER, space_after=0)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ================================================================
    # SLIDE 1 — PORTADA (fondo claro)
    # ================================================================
    s1 = prs.slides.add_slide(blank)

    _rect(s1, Inches(0), Inches(0), W, H, WHITE)
    _rect(s1, Inches(0), Inches(0), W, Inches(0.12), NAVY)
    _rect(s1, Inches(0), Inches(0.12), W, Inches(0.05), ORANGE)

    if LOGO_RETO.exists():
        s1.shapes.add_picture(str(LOGO_RETO), Inches(4.4), Inches(0.7), width=Inches(4.5))

    _rect(s1, Inches(2.0), Inches(3.7), Inches(9.3), Inches(0.04), ORANGE)

    tf1 = _text_box(s1, Inches(1.0), Inches(3.9), Inches(11.3), Inches(1.0))
    _add_text(tf1, "Aplicación de Monitoreo ReTo", 38, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    tf2 = _text_box(s1, Inches(1.0), Inches(4.85), Inches(11.3), Inches(0.6))
    _add_text(tf2, "Herramienta para detección, análisis y validación de discurso de odio", 18, color=GREY_TEXT, align=PP_ALIGN.CENTER)

    tf3 = _text_box(s1, Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.5))
    _add_text(tf3, "Presentación para socios del proyecto  ·  Reunión 7 de abril de 2026", 14, color=ORANGE, align=PP_ALIGN.CENTER)

    _rect(s1, Inches(0), H - Inches(1.3), W, Inches(1.3), GREY_BG)
    _rect(s1, Inches(0), H - Inches(1.32), W, Inches(0.04), ORANGE)

    x_logo = Inches(1.0)
    logo_h = Inches(0.55)
    y_logo = H - Inches(1.05)
    if LOGOS_DIR.exists():
        for png in sorted(LOGOS_DIR.glob("*.png")):
            try:
                s1.shapes.add_picture(str(png), x_logo, y_logo, height=logo_h)
                x_logo += Inches(1.8)
                if float(x_logo) > float(W) - Inches(1.5):
                    break
            except Exception:
                pass

    # ================================================================
    # SLIDE 2 — ¿QUÉ APORTA LA PLATAFORMA?
    # ================================================================
    s2 = prs.slides.add_slide(blank)
    _rect(s2, Inches(0), Inches(0), W, H, WHITE)
    _slide_title_bar(s2, "¿Qué aporta la plataforma hoy?", "Valor operativo para el proyecto")
    _add_footer_logos(s2)

    items = [
        ("Centralización", "Unifica información de mensajes en X (Twitter) y YouTube en un único entorno de trabajo."),
        ("Monitoreo", "Permite seguir volumen, tendencias, categorías y señales de interés por plataforma y medio."),
        ("Operatividad", "Facilita tareas de análisis, priorización y seguimiento del discurso de odio."),
        ("Validación humana", "Integra anotación y validación de mensajes para asegurar calidad y trazabilidad."),
        ("Uso inmediato", "Apoya tareas actuales del proyecto, más allá del proceso de mejora continua del modelo."),
    ]

    colors = [BLUE, NAVY, GREEN, PURPLE, ORANGE]

    for i, (title, desc) in enumerate(items):
        y = Inches(1.6) + Inches(i * 0.85)
        _rect(s2, Inches(0.8), y + Inches(0.05), Inches(0.08), Inches(0.6), colors[i])
        tf = _text_box(s2, Inches(1.1), y, Inches(10.5), Inches(0.8))
        p = tf.paragraphs[0]
        run_t = p.add_run()
        run_t.text = f"{title}:  "
        run_t.font.size = Pt(16)
        run_t.font.bold = True
        run_t.font.color.rgb = NAVY
        run_d = p.add_run()
        run_d.text = desc
        run_d.font.size = Pt(15)
        run_d.font.color.rgb = GREY_TEXT

    highlight = _rounded_rect(s2, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.6), ORANGE_LIGHT)
    htf = _text_box(s2, Inches(1.1), Inches(5.98), Inches(11.2), Inches(0.55))
    _add_text(htf, "La aplicación ya está operativa como herramienta de trabajo, en paralelo a la mejora continua del modelo.", 15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 3 — NAVEGACIÓN (PERFIL EDITOR)
    # ================================================================
    s3 = prs.slides.add_slide(blank)
    _rect(s3, Inches(0), Inches(0), W, H, WHITE)
    _slide_title_bar(s3, "Navegación de la aplicación", "Perfil Editor — acceso para socios")
    _add_footer_logos(s3)

    _rounded_rect(s3, Inches(0.6), Inches(1.6), Inches(4.2), Inches(1.4), RGBColor(0xE8, 0xF0, 0xFE))
    tf_acc = _text_box(s3, Inches(0.85), Inches(1.7), Inches(3.8), Inches(1.2))
    _add_text(tf_acc, "Acceso", 18, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    _add_text(tf_acc, "Ingreso con usuario y contraseña asignados a cada socio del proyecto.", 14, color=GREY_TEXT)

    _rounded_rect(s3, Inches(5.2), Inches(1.6), Inches(7.6), Inches(1.4), RGBColor(0xFD, 0xF0, 0xE0))
    tf_no = _text_box(s3, Inches(5.45), Inches(1.7), Inches(7.1), Inches(1.2))
    _add_text(tf_no, "Secciones no visibles para Editor", 18, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    _add_text(tf_no, "Comparativa modelos  ·  Calidad LLM (reservadas para administración).", 14, color=GREY_TEXT)

    sections = [
        ("Panel general", "Indicadores clave consolidados."),
        ("Categorías de odio", "Distribución por las 6 categorías ReTo."),
        ("Ranking de medios", "Top medios por volumen y % de odio."),
        ("Análisis contextual", "Tendencias semanales y alertas."),
        ("Términos frecuentes", "Palabras más recurrentes en mensajes candidatos."),
        ("Dataset Gold", "Evaluación del etiquetado validado."),
        ("Análisis Art. 510", "Potenciales delitos bajo el Código Penal."),
        ("Anotación y validación", "Flujo de trabajo de anotación humana."),
        ("Delitos de odio (oficial)", "Datos oficiales de España."),
    ]

    col_x = [Inches(0.6), Inches(4.8), Inches(9.0)]
    for i, (sec_name, sec_desc) in enumerate(sections):
        col = i % 3
        row = i // 3
        x = col_x[col]
        y = Inches(3.35) + Inches(row * 1.0)

        _rect(s3, x, y + Inches(0.03), Inches(0.06), Inches(0.7), BLUE)
        tf = _text_box(s3, x + Inches(0.2), y, Inches(3.7), Inches(0.85))
        _add_text(tf, sec_name, 14, bold=True, color=NAVY)
        _add_text(tf, sec_desc, 12, color=GREY_TEXT, space_after=0)

    # ================================================================
    # SLIDE 4 — FLUJO DE USO
    # ================================================================
    s4 = prs.slides.add_slide(blank)
    _rect(s4, Inches(0), Inches(0), W, H, GREY_BG)
    _slide_title_bar(s4, "Flujo de uso para socios", "Cuatro pasos del trabajo diario")
    _add_footer_logos(s4)

    steps = [
        (1, "Monitorear", "Revisar paneles y métricas\npor plataforma, categoría\ny medio de comunicación.", BLUE),
        (2, "Explorar", "Identificar patrones,\npicos contextuales y\nfocos de interés.", GREEN),
        (3, "Validar", "Anotar y validar mensajes\nen la sección\ncorrespondiente.", ORANGE),
        (4, "Retroalimentar", "La validación humana\nmejora la calidad operativa\ny el modelo.", PURPLE),
    ]

    start_x = Inches(0.55)
    for i, (num, title, desc, color) in enumerate(steps):
        _step_card(s4, start_x + Inches(i * 3.15), Inches(1.75), num, title, desc, color)

    for i in range(3):
        arrow_x = start_x + Inches((i + 1) * 3.15) - Inches(0.35)
        arrow = s4.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(3.25), Inches(0.5), Inches(0.35)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        arrow.line.fill.background()

    note_tf = _text_box(s4, Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.5))
    _add_text(note_tf, "La validación humana fortalece la utilidad operativa inmediata y mejora progresivamente el desempeño del sistema.", 14, bold=False, color=GREY_TEXT, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 5 — ANTES DE LA REUNIÓN
    # ================================================================
    s5 = prs.slides.add_slide(blank)
    _rect(s5, Inches(0), Inches(0), W, H, WHITE)
    _slide_title_bar(s5, "Antes de la reunión del 7 de abril", "Interacción previa y propuestas de mejora")
    _add_footer_logos(s5)

    intro_tf = _text_box(s5, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.8))
    _add_text(intro_tf, "Es clave que cada socio interactúe con la aplicación antes de la reunión para traer observaciones y propuestas concretas.", 16, color=DARK_TEXT, align=PP_ALIGN.LEFT)

    checks = [
        "Revisar al menos tres secciones de análisis (Panel general, Ranking de medios, Categorías...).",
        "Probar la sección de Anotación y validación con algunos mensajes.",
        "Identificar mejoras de usabilidad y navegación.",
        "Evaluar qué información adicional sería útil visualizar.",
        "Pensar en la utilidad concreta para las tareas de cada entidad.",
    ]

    check_y = Inches(2.5)
    for i, text in enumerate(checks):
        y = check_y + Inches(i * 0.65)
        check_circ = _circle(s5, Inches(1.0), y + Inches(0.05), Inches(0.35), BLUE)
        ctf = check_circ.text_frame
        ctf.paragraphs[0].text = "\u2713"
        ctf.paragraphs[0].font.size = Pt(16)
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

        tf = _text_box(s5, Inches(1.55), y, Inches(10.8), Inches(0.5))
        _add_text(tf, text, 15, color=DARK_TEXT)

    highlight2 = _rounded_rect(s5, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.65), NAVY)
    htf2 = _text_box(s5, Inches(1.1), Inches(5.95), Inches(11.2), Inches(0.55))
    _add_text(htf2, "Objetivo: consolidar la plataforma como herramienta de trabajo compartida del proyecto ReTo.", 17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ================================================================
    out = SCRIPT_DIR / OUT_NAME
    prs.save(str(out))
    print(f"Generado: {out}")
    return out


if __name__ == "__main__":
    build()
